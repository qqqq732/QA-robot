import os
import tempfile
import traceback
import shutil
import asyncio  
import numpy as np  
from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List, Optional, Tuple
from fastapi.responses import StreamingResponse

# 向量库与切片相关组件
from langchain_community.vectorstores import FAISS
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

# 本地多格式文件解析库
import pdfplumber
import pypdfium2 as pdfium
import easyocr  
from docx import Document as DocxDocument
import openpyxl
from pptx import Presentation  

# 导入百炼原生 SDK 组件
import dashscope
from dashscope import TextEmbedding
from dashscope import Generation

# ==========================================
# 基础配置
# ==========================================
DASHSCOPE_API_KEY = "sk-密钥"
os.makedirs("vector_store", exist_ok=True)
dashscope.api_key = DASHSCOPE_API_KEY


# ==========================================
# 自定义 Embedding 类
# ==========================================
class DashScopeEmbeddings:
    def __init__(self, model: str = "text-embedding-v3"):
        self.model = model

    def embed_query(self, text: str) -> List[float]:
        try:
            resp = TextEmbedding.call(
                model=self.model,
                input=text,
                text_type="document"
            )
            if resp.status_code == 200:
                return resp.output['embeddings'][0]['embedding']
            else:
                raise Exception(f"Embedding API 失败: {resp.message}")
        except Exception as e:
            print(f"Embedding 错误: {e}")
            raise

    def embed_documents(self, texts: List[str]) -> List[List[float]]:
        return [self.embed_query(text) for text in texts]

    def __call__(self, text: str) -> List[float]:
        return self.embed_query(text)


# 初始化向量实例
embedding = DashScopeEmbeddings(model="text-embedding-v3")

app = FastAPI(title="RAG文档问答后端")

# 跨域配置
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ==========================================
# 入参模型
# ==========================================
class Message(BaseModel):
    role: str
    content: str


class ChatRequest(BaseModel):
    query: str
    history: List[Message] = []
    kb_id: Optional[str] = "default"
    model_name: Optional[str] = "qwen-max"


print("正在初始化 EasyOCR 引擎...")
ocr_engine = easyocr.Reader(['ch_sim', 'en'], gpu=False)


# ==========================================
# 文件解析、文本切块
# ==========================================
def extract_text_from_file(file_path: str, ext: str) -> str:
    text = ""
    ext = ext.lower()
    try:
        if ext in [".txt", ".md"]:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()

        elif ext == ".pdf":
            # 1. 先尝试用 pdfplumber 提取电子文本
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
            except Exception as pdf_err:
                print(f"pdfplumber 常规解析失败: {pdf_err}")

            if len(text.strip()) < 10:
                print("检测到该 PDF 可能是图片扫描件，启动备用 EasyOCR 引擎进行识别...")
                text = ""

                with pdfium.PdfDocument(file_path) as pdf_render:
                    for i, page in enumerate(pdf_render):
                        print(f"--- 正在使用 EasyOCR 识别第 {i + 1} 页 ---")
                        bitmap = page.render(scale=2)  # 放大2倍确保清晰度
                        pil_img = bitmap.to_pil().convert("RGB")
                        img_np = np.array(pil_img)

                        # 使用 EasyOCR 提取文本
                        result = ocr_engine.readtext(img_np, detail=0)
                        if result:
                            page_text = " ".join(result)
                            text += page_text + "\n"
                            print(f"第 {i + 1} 页识别成功，提取到 {len(page_text)} 个字。")
                        else:
                            print(f"⚠️ 警告：第 {i + 1} 页未识别到任何文字。")

                print(f"💡 本地 OCR 识别成功完成！总字数: {len(text)}")

        elif ext == ".docx":
            doc = DocxDocument(file_path)
            for para in doc.paragraphs:
                if para.text: text += para.text + "\n"

        elif ext == ".xlsx":
            wb = openpyxl.load_workbook(file_path)
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip(): text += row_text + "\n"

        elif ext in [".ppt", ".pptx"]:
            print(f"正在解析 PPT 文件: {file_path}")
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                if slide_text:
                    text += f"[第 {i + 1} 页幻灯片]\n" + "\n".join(slide_text) + "\n\n"
            print(f"💡 PPT 解析成功！总字数: {len(text)}")

    except Exception as e:
        print(f"解析文件错误: {e}")
        traceback.print_exc()

    return text.strip()


def create_chunks(text: str):
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    return splitter.split_text(text)


def retrieve_context_with_sources(query: str, kb_id: str) -> Tuple[str, List[str]]:
    try:
        db_path = f"vector_store/{kb_id}"
        if not os.path.exists(f"{db_path}/index.faiss"):
            return "", []
        vectorstore = FAISS.load_local(db_path, embedding, allow_dangerous_deserialization=True)
        docs = vectorstore.similarity_search(query, k=3)

        context_str = "\n\n".join([d.page_content for d in docs])
        sources = list(set([d.metadata.get("source", "未知文档") for d in docs if d.metadata]))
        return context_str, sources
    except Exception as e:
        print(f"检索异常: {e}")
        return "", []


# ==========================================
# 接口实现
# ==========================================
@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...), kb_id: str = Form("default")):
    tmp_path = None
    try:
        ext = os.path.splitext(file.filename)[-1].lower()
        allowed = [".txt", ".md", ".pdf", ".docx", ".xlsx", ".ppt", ".pptx"]
        if ext not in allowed:
            return {"status": "error", "msg": f"不支持的文件格式: {ext}"}

        content = await file.read()
        if not content:
            return {"status": "error", "msg": "文件内容为空"}

        with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
            tmp.write(content)
            tmp_path = tmp.name

        text = await asyncio.to_thread(extract_text_from_file, tmp_path, ext)

        if not text or len(text) < 10:
            return {"status": "error", "msg": "无法从文件中提取出足够的有效文本内容"}

        chunks = create_chunks(text)
        documents = [Document(page_content=c, metadata={"source": file.filename, "chunk_id": i})
                     for i, c in enumerate(chunks)]

        db_path = f"vector_store/{kb_id}"
        if os.path.exists(f"{db_path}/index.faiss"):
            vectorstore = FAISS.load_local(db_path, embedding, allow_dangerous_deserialization=True)
            vectorstore.add_documents(documents)
        else:
            vectorstore = FAISS.from_documents(documents, embedding)

        vectorstore.save_local(db_path)

        filename_file = os.path.join(db_path, "filenames.txt")
        existing_files = set()
        if os.path.exists(filename_file):
            with open(filename_file, "r", encoding="utf-8") as f:
                existing_files = {line.strip() for line in f if line.strip()}
        existing_files.add(file.filename)
        with open(filename_file, "w", encoding="utf-8") as f:
            f.write("\n".join(existing_files))

        return {"status": "ok", "msg": "上传成功"}
    except Exception as e:
        return {"status": "error", "msg": str(e)}
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)


@app.post("/api/chat/stream")
async def chat_stream(request: ChatRequest):
    async def generate():
        try:
            current_model = request.model_name
            if "-vl-" in current_model.lower():
                current_model = "qwen-plus"

            context, sources = retrieve_context_with_sources(request.query, request.kb_id)

            if context and sources:
                source_titles = "、".join([f"《{s}》" for s in sources])
                system_prompt = f"""你是一个极其严谨且排版规范的文档问答机器人。
【参考资料】
{context}

【核心行为准则 - 必须严格执行】
1. **只能**基于上述提供的参考资料进行回答。
2. 绝对不能加入任何参考资料中没有提及的个人推测、联想或常识扩展。
3. 如果参考资料中没有包含回答用户问题所需的信息，请直接回答：“未找到相关内容。”，严禁胡编乱造。
4. 回答的开头**必须**带有固定引用声明，格式为：“基于{source_titles}文档回答：\n\n”

【排版与断行强制要求】
5. 必须使用标准的 Markdown 有序列表（1., 2., 3.）进行分点回答。
6. 每个要点之间，必须使用两个换行符进行强制隔开！绝对禁止把多个要点挤在同一行。
7. 每一点的标题和后面的具体描述内容之间，必须换行，或者使用冒号加空格。
"""
            else:
                system_prompt = """你是一个严谨的文档问答机器人。
【重要提示】当前在本地知识库中未找到与用户问题相关的参考资料。
请严格按照以下要求回答：
1. 先在回答的第一行明确提示用户：“（提示：未在上传的知识库文档中找到相关内容，以下回答基于 AI 通用知识库：）\n\n”
2. 然后再根据你原本掌握的知识，规范地回答用户的问题。
3. 每个要点之间，必须使用两个换行符强制分段！必须使用标准 Markdown 格式并带有清晰的数字序号（如 1., 2., 3.）。
"""

            messages = [{"role": "system", "content": system_prompt}]
            for m in request.history:
                messages.append({"role": m.role, "content": m.content})
            messages.append({"role": "user", "content": request.query})

            # 💡 核心优化：百炼的 SDK 内部是同步阻塞 IO，使用 asyncio.to_thread 包裹使其兼容异步生成器
            responses = await asyncio.to_thread(
                Generation.call,
                model=current_model,
                messages=messages,
                result_format='message',
                stream=True,
                incremental_output=True
            )

            for response in responses:
                if response.status_code == 200:
                    choice = response.output.choices[0]
                    message = getattr(choice, 'message', {})
                    reasoning, content = "", ""
                    try:
                        reasoning = getattr(message, 'reasoning_content', '')
                    except:
                        pass
                    try:
                        content = getattr(message, 'content', '')
                    except:
                        pass

                    if reasoning:
                        safe_reasoning = reasoning.replace("\n", "[BR]")
                        yield f"data: <think>{safe_reasoning}\n\n"

                    if content:
                        safe_content = content.replace("\n", "[BR]")
                        yield f"data: {safe_content}\n\n"

                    await asyncio.sleep(0.01)
                else:
                    yield f"data: 错误({response.code}): {response.message}\n\n"
                    break
            yield "data: [DONE]\n\n"
        except Exception as e:
            traceback.print_exc()
            yield f"data: 错误: {str(e)}\n\n"

    return StreamingResponse(
        generate(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"
        }
    )


@app.get("/api/kb/list")
async def list_kb():
    kbs = []
    if os.path.exists("vector_store"):
        for d in os.listdir("vector_store"):
            dir_path = os.path.join("vector_store", d)
            if os.path.isdir(dir_path) and os.path.exists(os.path.join(dir_path, "index.faiss")):
                kbs.append(d)
    return {"knowledge_bases": kbs}


@app.delete("/api/kb/{kb_id}")
async def delete_kb(kb_id: str):
    kb_folder = os.path.join("vector_store", kb_id)
    if os.path.isdir(kb_folder):
        shutil.rmtree(kb_folder)
        return {"code": 200, "msg": f"知识库【{kb_id}】已成功删除"}
    return {"code": 400, "msg": f"知识库【{kb_id}】不存在"}


@app.delete("/api/session/{session_id}")
async def delete_session_kb(session_id: str):
    kb_path = os.path.join("vector_store", session_id)
    if os.path.exists(kb_path):
        shutil.rmtree(kb_path)
        return {"status": "ok", "msg": "会话知识库已注销"}
    return {"status": "ok", "msg": "无关联库"}


@app.get("/api/kb/files/{kb_id}")
async def get_kb_files(kb_id: str):
    try:
        fn = os.path.join("vector_store", kb_id, "filenames.txt")
        if os.path.exists(fn):
            with open(fn, "r", encoding="utf-8") as f:
                lines = [l.strip() for l in f if l.strip()]
            return {"kb_id": kb_id, "filenames": lines}
    except:
        pass
    return {"kb_id": kb_id, "filenames": []}


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
